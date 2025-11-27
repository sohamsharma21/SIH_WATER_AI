"""
SIH Presentation Generator
Auto-generates SIH presentation using python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime
from pathlib import Path

# Output path
OUTPUT_DIR = Path(__file__).parent.parent / "docs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT_FILE = OUTPUT_DIR / "SIH_presentation.pptx"


def create_presentation():
    """Create SIH presentation with all slides."""
    prs = Presentation()
    
    # Set slide size to 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Title Slide
    create_title_slide(prs)
    
    # Problem Statement
    create_problem_slide(prs)
    
    # Solution Overview
    create_solution_slide(prs)
    
    # Architecture
    create_architecture_slide(prs)
    
    # ML Models
    create_ml_models_slide(prs)
    
    # Digital Twin
    create_digital_twin_slide(prs)
    
    # Treatment Optimization
    create_treatment_optimization_slide(prs)
    
    # Results
    create_results_slide(prs)
    
    # Technology Stack
    create_technology_stack_slide(prs)
    
    # Future Scope
    create_future_scope_slide(prs)
    
    # Team
    create_team_slide(prs)
    
    # Thank You
    create_thank_you_slide(prs)
    
    # Save presentation
    prs.save(OUTPUT_FILE)
    print(f"✓ Presentation generated: {OUTPUT_FILE}")
    
    return OUTPUT_FILE


def create_title_slide(prs):
    """Create title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "SIH WATER AI"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(30, 58, 138)  # Blue
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_top = Inches(3.2)
    subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Industrial Wastewater Treatment Optimization System"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(59, 130, 246)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Team
    team_top = Inches(4.5)
    team_box = slide.shapes.add_textbox(left, team_top, width, Inches(0.5))
    team_frame = team_box.text_frame
    team_frame.text = "Team: Nova_Minds"
    team_para = team_frame.paragraphs[0]
    team_para.font.size = Pt(18)
    team_para.font.color.rgb = RGBColor(107, 114, 128)
    team_para.alignment = PP_ALIGN.CENTER


def create_problem_slide(prs):
    """Create problem statement slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    
    # Title
    title = slide.shapes.title
    title.text = "Problem Statement"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 58, 138)
    
    # Content
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Industrial Wastewater Challenges"
    
    p = tf.add_paragraph()
    p.text = "• High contamination levels affecting treatment efficiency"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "• Manual monitoring and optimization processes"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "• Lack of real-time insights and predictive capabilities"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "• Inefficient resource usage (energy, chemicals, time)"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "• Need for automated treatment optimization"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "• Limited visibility into plant operations"
    p.level = 0
    p.font.size = Pt(20)


def create_solution_slide(prs):
    """Create solution overview slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Solution Overview"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 58, 138)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "AI-Powered Treatment Optimization Platform"
    
    p = tf.add_paragraph()
    p.text = "✓ Multi-Model AI: 4 trained ML models for accurate predictions"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "✓ Digital Twin: 3D real-time visualization of treatment plant"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "✓ Real-time Monitoring: MQTT sensor ingestion with instant alerts"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "✓ Automated Optimization: AI-driven treatment recommendations"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "✓ Comprehensive Reports: PDF generation with analysis"
    p.level = 0
    p.font.size = Pt(20)


def create_architecture_slide(prs):
    """Create architecture slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "System Architecture"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 58, 138)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Full-Stack Architecture"
    
    p = tf.add_paragraph()
    p.text = "Frontend (Next.js) → Backend (FastAPI) → Database (Supabase)"
    p.level = 0
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "MQTT Broker → Real-time Sensor Data → ML Models → Predictions"
    p.level = 0
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Digital Twin → 3D Visualization → Treatment Optimization → Reports"
    p.level = 0
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Key Components:"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• React-Three-Fiber for 3D visualization"
    p.level = 0
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• scikit-learn for ML models"
    p.level = 0
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• ReportLab for PDF generation"
    p.level = 0
    p.font.size = Pt(18)


def create_ml_models_slide(prs):
    """Create ML models slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Multi-Model AI System"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 58, 138)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "4 Trained ML Models"
    
    p = tf.add_paragraph()
    p.text = "1. NYC DEP Model: Primary/Secondary treatment prediction"
    p.level = 0
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "2. Water Potability Model: Tertiary treatment classification"
    p.level = 0
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "3. UCI Model: Contamination severity assessment"
    p.level = 0
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "4. Full-Scale WWTP Model: Aeration control & BOD/COD prediction"
    p.level = 0
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Unified Pipeline: Imputer → PolynomialFeatures → StandardScaler → RandomForest"
    p.level = 0
    p.font.size = Pt(16)
    p.font.italic = True


def create_digital_twin_slide(prs):
    """Create digital twin slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "3D Digital Twin"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 58, 138)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Real-time 3D Visualization"
    
    p = tf.add_paragraph()
    p.text = "✓ 3D plant model with tanks, pipes, clarifiers"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "✓ Real-time water level animations"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "✓ Color changes based on turbidity/contamination"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "✓ Aeration bubble effects"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "✓ Camera controls for navigation"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "✓ Hover overlays for parameter display"
    p.level = 0
    p.font.size = Pt(20)


def create_treatment_optimization_slide(prs):
    """Create treatment optimization slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Treatment Optimization"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 58, 138)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Three-Stage AI-Driven Optimization"
    
    p = tf.add_paragraph()
    p.text = "Primary: Settling time, Coagulant dosing, Sludge volume index"
    p.level = 0
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Secondary: Aeration time, DO control, Blower speed, Sludge age"
    p.level = 0
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Tertiary: Filtration rate, Chlorine dosing, RO trigger"
    p.level = 0
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Final Reuse: Irrigation/Industrial/Environmental/Drinking classification"
    p.level = 0
    p.font.size = Pt(18)


def create_results_slide(prs):
    """Create results slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Results & Impact"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 58, 138)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Achievements"
    
    p = tf.add_paragraph()
    p.text = "✓ 4 ML models trained on real wastewater datasets"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "✓ Real-time sensor data ingestion via MQTT"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "✓ 3D digital twin with live visualizations"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "✓ Automated treatment recommendations"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "✓ Comprehensive PDF report generation"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "✓ Production-ready full-stack platform"
    p.level = 0
    p.font.size = Pt(20)


def create_technology_stack_slide(prs):
    """Create technology stack slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Technology Stack"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 58, 138)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Modern Tech Stack"
    
    p = tf.add_paragraph()
    p.text = "Frontend: Next.js 14, React, TypeScript, TailwindCSS, Three.js"
    p.level = 0
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Backend: FastAPI, Python 3.11, Supabase, scikit-learn"
    p.level = 0
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Database: PostgreSQL (Supabase), Row Level Security"
    p.level = 0
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "ML: scikit-learn, pandas, numpy, joblib"
    p.level = 0
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "3D: React-Three-Fiber, Drei, GSAP"
    p.level = 0
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "IoT: MQTT (paho-mqtt), Real-time subscriptions"
    p.level = 0
    p.font.size = Pt(18)


def create_future_scope_slide(prs):
    """Create future scope slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Future Scope"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 58, 138)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Enhancement Roadmap"
    
    p = tf.add_paragraph()
    p.text = "• Real-time WebSocket updates for digital twin"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "• Advanced visualization with Plotly 3D graphs"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "• Mobile app support"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "• Historical data analysis and trending"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "• Automated alert system"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "• Integration with SCADA systems"
    p.level = 0
    p.font.size = Pt(20)


def create_team_slide(prs):
    """Create team slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Team"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(30, 58, 138)
    title_para.alignment = PP_ALIGN.CENTER
    
    team_top = Inches(3.5)
    team_box = slide.shapes.add_textbox(left, team_top, width, Inches(1))
    team_frame = team_box.text_frame
    team_frame.text = "Nova_Minds"
    team_para = team_frame.paragraphs[0]
    team_para.font.size = Pt(36)
    team_para.font.color.rgb = RGBColor(59, 130, 246)
    team_para.alignment = PP_ALIGN.CENTER


def create_thank_you_slide(prs):
    """Create thank you slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Thank You"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(30, 58, 138)
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle_top = Inches(3.5)
    subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Questions?"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(107, 114, 128)
    subtitle_para.alignment = PP_ALIGN.CENTER


if __name__ == "__main__":
    print("=" * 60)
    print("Generating SIH Presentation...")
    print("=" * 60)
    
    output_file = create_presentation()
    
    print("=" * 60)
    print(f"✓ Presentation saved to: {output_file}")
    print("=" * 60)
    print("\nNote: Add screenshots and custom images manually to slides as needed.")

